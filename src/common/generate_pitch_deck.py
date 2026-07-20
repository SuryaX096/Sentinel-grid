import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_pitch_deck():
    prs = Presentation()
    # Set widescreen layout (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom Theme Colors
    BG_COLOR = RGBColor(11, 15, 25)         # Dark Navy (#0b0f19)
    TITLE_COLOR = RGBColor(56, 189, 248)    # Cyan / Sky Blue (#38bdf8)
    TEXT_COLOR = RGBColor(226, 232, 240)     # Light Gray (#e2e8f0)
    ACCENT_COLOR = RGBColor(245, 158, 11)   # Orange (#f59e0b)
    MUTED_COLOR = RGBColor(148, 163, 184)    # Slate Gray (#94a3b8)
    CARD_BG = RGBColor(30, 41, 59)          # Card Navy (#1e293b)
    CARD_BORDER = RGBColor(51, 65, 85)      # Card Border (#334155)

    blank_slide_layout = prs.slide_layouts[6] # Blank slide layout

    def apply_dark_background(slide):
        # Draw a full-slide rectangle as background
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG_COLOR
        rect.line.fill.background() # No border
        # Send to back is done implicitly when added first

    def add_slide_header(slide, title_text, category_text="SENTINELMIND AI"):
        # Header category text
        txBox_cat = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.3))
        tf_cat = txBox_cat.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_COLOR
        p_cat.font.name = "Outfit"

        # Main slide title
        txBox_title = slide.shapes.add_textbox(Inches(0.75), Inches(0.75), Inches(11.83), Inches(0.8))
        tf_title = txBox_title.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = TITLE_COLOR
        p_title.font.name = "Outfit"

    # ================= Slide 1: Title Slide =================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide1)

    # Decorative visual bar
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.2), Inches(3.5), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    # Title text box
    title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.83), Inches(3.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0

    p1 = tf1.paragraphs[0]
    p1.text = "SentinelMind AI"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = TITLE_COLOR
    p1.font.name = "Outfit"

    p2 = tf1.add_paragraph()
    p2.text = "Autonomous Multi-Agent Threat Detection, Attribution, & Response"
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_COLOR
    p2.font.name = "Outfit"
    p2.space_before = Pt(15)

    p3 = tf1.add_paragraph()
    p3.text = "ET AI Hackathon 2026  |  Problem Statement 7  |  Submission Pitch Deck"
    p3.font.size = Pt(14)
    p3.font.color.rgb = MUTED_COLOR
    p3.font.name = "Outfit"
    p3.space_before = Pt(45)


    # ================= Slide 2: The Core Problem =================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide2)
    add_slide_header(slide2, "The Cyber Security Challenge", "The Opportunity")

    # Column 1: The Problem
    col1_box = slide2.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(5.6), Inches(4.5))
    tf2_col1 = col1_box.text_frame
    tf2_col1.word_wrap = True
    tf2_col1.margin_left = tf2_col1.margin_right = tf2_col1.margin_top = tf2_col1.margin_bottom = 0

    p_c1_title = tf2_col1.paragraphs[0]
    p_c1_title.text = "❌ The Industry Status Quo"
    p_c1_title.font.size = Pt(20)
    p_c1_title.font.bold = True
    p_c1_title.font.color.rgb = ACCENT_COLOR
    p_c1_title.space_after = Pt(15)

    problems = [
        "Static Rules Fall Short: Legacy firewalls and SIEM systems rely on signature matches, letting zero-day exploits slide past.",
        "Attribution Blindspots: Detections lack context. Security teams receive alerts but cannot match them to actual adversary techniques.",
        "Response Starvation: Alert fatigue delays responses. Hand-crafted actions are slow, while un-gated scripts risk breaking business systems."
    ]
    for prob in problems:
        p_c1_item = tf2_col1.add_paragraph()
        p_c1_item.text = prob
        p_c1_item.font.size = Pt(14)
        p_c1_item.font.color.rgb = TEXT_COLOR
        p_c1_item.space_after = Pt(12)
        p_c1_item.level = 0

    # Column 2: The Solution
    col2_box = slide2.shapes.add_textbox(Inches(6.98), Inches(2.0), Inches(5.6), Inches(4.5))
    tf2_col2 = col2_box.text_frame
    tf2_col2.word_wrap = True
    tf2_col2.margin_left = tf2_col2.margin_right = tf2_col2.margin_top = tf2_col2.margin_bottom = 0

    p_c2_title = tf2_col2.paragraphs[0]
    p_c2_title.text = "🛡️ The SentinelMind AI Solution"
    p_c2_title.font.size = Pt(20)
    p_c2_title.font.bold = True
    p_c2_title.font.color.rgb = TITLE_COLOR
    p_c2_title.space_after = Pt(15)

    solutions = [
        "Unsupervised Anomaly Detection: Learns 'normal' baseline network behavior to isolate unknown traffic profiles.",
        "Semantic Threat Intelligence Mapping: Automatically translates raw connection metrics into named MITRE ATT&CK techniques.",
        "Automated Gated Action: Triggers instant, blast-radius-aware mitigations, asking for human approval before high-impact changes."
    ]
    for sol in solutions:
        p_c2_item = tf2_col2.add_paragraph()
        p_c2_item.text = sol
        p_c2_item.font.size = Pt(14)
        p_c2_item.font.color.rgb = TEXT_COLOR
        p_c2_item.space_after = Pt(12)
        p_c2_item.level = 0


    # ================= Slide 3: Three-Agent Architecture =================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide3)
    add_slide_header(slide3, "Autonomous Three-Agent Pipeline", "System Architecture")

    # Three horizontal columns representing the agents
    width_col = Inches(3.6)
    spacing_col = Inches(0.4)
    start_left = Inches(0.75)
    top_y = Inches(2.0)
    height_y = Inches(4.5)

    agents = [
        {
            "num": "01",
            "name": "HERO AGENT",
            "subtitle": "Anomaly Detection Engine",
            "tech": "Isolation Forest Model",
            "desc": [
                "Ingests raw network connection metrics.",
                "Scores connection anomaly based on outliers.",
                "Computes Z-scores for every connection attribute to flag specific network indicator deviations."
            ],
            "accent": RGBColor(56, 189, 248)
        },
        {
            "num": "02",
            "name": "ATTRIBUTION AGENT",
            "subtitle": "Threat Semantic Classifier",
            "tech": "ChromaDB & Gemini API",
            "desc": [
                "Translates raw anomalous attributes to clear security text descriptions.",
                "Matches descriptions semantically with ChromaDB vector index of MITRE ATT&CK techniques.",
                "Optionally uses Gemini 1.5 Flash to validate candidates."
            ],
            "accent": ACCENT_COLOR
        },
        {
            "num": "03",
            "name": "SUPPORTING AGENT",
            "subtitle": "Playbook Orchestrator",
            "tech": "Decision State Machine",
            "desc": [
                "Looks up target asset details, owners, and active vulnerability profiles (CVEs).",
                "Applies risk scoring: risk = (anomaly score * 6.0) + (asset criticality weight).",
                "Triggers auto-execution or gates actions for manual SOC analyst approval."
            ],
            "accent": RGBColor(34, 197, 94)
        }
    ]

    for idx, agent in enumerate(agents):
        x_left = start_left + idx * (width_col + spacing_col)
        
        # Background card
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_left, top_y, width_col, height_y)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.5)

        # Content text box
        tb = slide3.shapes.add_textbox(x_left + Inches(0.25), top_y + Inches(0.25), width_col - Inches(0.5), height_y - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        p_num = tf.paragraphs[0]
        p_num.text = agent["num"]
        p_num.font.size = Pt(36)
        p_num.font.bold = True
        p_num.font.color.rgb = agent["accent"]

        p_name = tf.add_paragraph()
        p_name.text = agent["name"]
        p_name.font.size = Pt(18)
        p_name.font.bold = True
        p_name.font.color.rgb = TITLE_COLOR
        p_name.space_before = Pt(8)

        p_sub = tf.add_paragraph()
        p_sub.text = agent["subtitle"]
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = MUTED_COLOR

        p_tech = tf.add_paragraph()
        p_tech.text = f"Tech Stack: {agent['tech']}"
        p_tech.font.size = Pt(12)
        p_tech.font.bold = True
        p_tech.font.color.rgb = agent["accent"]
        p_tech.space_before = Pt(8)
        p_tech.space_after = Pt(12)

        for bullet in agent["desc"]:
            p_b = tf.add_paragraph()
            p_b.text = "• " + bullet
            p_b.font.size = Pt(12)
            p_b.font.color.rgb = TEXT_COLOR
            p_b.space_after = Pt(8)


    # ================= Slide 4: Hero Agent - Deep Dive =================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide4)
    add_slide_header(slide4, "Hero Agent: Anomaly Detection", "Machine Learning Core")

    # Left: Details
    col1_box = slide4.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(5.6), Inches(4.5))
    tf4_col1 = col1_box.text_frame
    tf4_col1.word_wrap = True
    tf4_col1.margin_left = tf4_col1.margin_right = tf4_col1.margin_top = tf4_col1.margin_bottom = 0

    p_h = tf4_col1.paragraphs[0]
    p_h.text = "How the Model Operates"
    p_h.font.size = Pt(20)
    p_h.font.bold = True
    p_h.font.color.rgb = TITLE_COLOR
    p_h.space_after = Pt(15)

    bullets = [
        "Unsupervised Approach: Trained on raw connection flows from the standard NSL-KDD dataset, filtering to baseline clean/normal records only.",
        "Feature Engineering: Real-time calculation of logarithmic duration and size, outbound/inbound ratios, and connection service frequency.",
        "Deviation Profiling: Calculates Z-score outliers against the training population to generate detailed security indicator tags."
    ]
    for bullet in bullets:
        p_b = tf4_col1.add_paragraph()
        p_b.text = bullet
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_COLOR
        p_b.space_after = Pt(15)

    # Right: Metrics Card
    x_card = Inches(6.98)
    card_right = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_card, Inches(2.0), Inches(5.6), Inches(4.5))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = CARD_BG
    card_right.line.color.rgb = CARD_BORDER
    card_right.line.width = Pt(1.5)

    tb_metrics = slide4.shapes.add_textbox(x_card + Inches(0.4), Inches(2.4), Inches(4.8), Inches(3.7))
    tf_m = tb_metrics.text_frame
    tf_m.word_wrap = True
    tf_m.margin_left = tf_m.margin_right = tf_m.margin_top = tf_m.margin_bottom = 0

    p_mt = tf_m.paragraphs[0]
    p_mt.text = "Holdout Set Evaluation Metrics"
    p_mt.font.size = Pt(18)
    p_mt.font.bold = True
    p_mt.font.color.rgb = ACCENT_COLOR
    p_mt.space_after = Pt(15)

    metrics = [
        ("Precision", "95.82%", "Low false alarm rate on clean connections"),
        ("Recall (Sensitivity)", "74.47%", "Successfully detects hidden malicious flows"),
        ("F1-Score", "0.8381", "High harmonic average performance profile"),
        ("False Positive Rate", "4.29%", "Less than 5% normal traffic flagged")
    ]
    for label, val, desc in metrics:
        p_m = tf_m.add_paragraph()
        p_m.text = f"{label}: "
        p_m.font.size = Pt(14)
        p_m.font.bold = True
        p_m.font.color.rgb = TEXT_COLOR
        
        # Add value with cyan highlight
        r = p_m.add_run()
        r.text = val
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = TITLE_COLOR

        # Subtext description
        p_d = tf_m.add_paragraph()
        p_d.text = f"   ({desc})"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = MUTED_COLOR
        p_d.space_after = Pt(8)


    # ================= Slide 5: Attribution Agent - Deep Dive =================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide5)
    add_slide_header(slide5, "Attribution Agent: Threat Mapping", "Threat Intelligence")

    # Left: How it works
    col1_box = slide5.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(5.6), Inches(4.5))
    tf5_col1 = col1_box.text_frame
    tf5_col1.word_wrap = True
    tf5_col1.margin_left = tf5_col1.margin_right = tf5_col1.margin_top = tf5_col1.margin_bottom = 0

    p_h5 = tf5_col1.paragraphs[0]
    p_h5.text = "Bridging ML Outliers to Adversary Context"
    p_h5.font.size = Pt(20)
    p_h5.font.bold = True
    p_h5.font.color.rgb = TITLE_COLOR
    p_h5.space_after = Pt(15)

    bullets5 = [
        "Feature Translation: Maps flagged indicators (e.g. serror_rate, wrong_fragment) to natural English descriptions of threat behaviors.",
        "ChromaDB Indexing: Embedded all primary MITRE ATT&CK techniques using 'all-MiniLM-L6-v2' (generating 384-dimension vector embeddings).",
        "Semantic Search: Queries the database with the threat text profile to find the top 3 closest matching ATT&CK candidates."
    ]
    for bullet in bullets5:
        p_b = tf5_col1.add_paragraph()
        p_b.text = bullet
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_COLOR
        p_b.space_after = Pt(15)

    # Right: LLM Validation
    card_r5 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(2.0), Inches(5.6), Inches(4.5))
    card_r5.fill.solid()
    card_r5.fill.fore_color.rgb = CARD_BG
    card_r5.line.color.rgb = CARD_BORDER
    card_r5.line.width = Pt(1.5)

    tb_val = slide5.shapes.add_textbox(Inches(6.98) + Inches(0.4), Inches(2.4), Inches(4.8), Inches(3.7))
    tf_v = tb_val.text_frame
    tf_v.word_wrap = True
    tf_v.margin_left = tf_v.margin_right = tf_v.margin_top = tf_v.margin_bottom = 0

    p_vtitle = tf_v.paragraphs[0]
    p_vtitle.text = "Dual-Mode Resolution Engine"
    p_vtitle.font.size = Pt(18)
    p_vtitle.font.bold = True
    p_vtitle.font.color.rgb = ACCENT_COLOR
    p_vtitle.space_after = Pt(15)

    modes = [
        ("1. High-Fidelity LLM Mode (API Connected)", 
         "Validates search candidates using Gemini 1.5 Flash. Cross-references indicators and outputs a JSON object with final technique selection and a detailed context rationale."),
        ("2. Local Similarity Mode (Offline Fallback)", 
         "If no API keys are set, attributes automatically to the top semantic vector search candidate. This maintains 100% functionality without internet connections.")
    ]
    for m_title, m_desc in modes:
        p_mt = tf_v.add_paragraph()
        p_mt.text = m_title
        p_mt.font.size = Pt(14)
        p_mt.font.bold = True
        p_mt.font.color.rgb = TEXT_COLOR
        
        p_md = tf_v.add_paragraph()
        p_md.text = m_desc
        p_md.font.size = Pt(12)
        p_md.font.color.rgb = MUTED_COLOR
        p_md.space_after = Pt(12)


    # ================= Slide 6: Supporting Agent - Deep Dive =================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide6)
    add_slide_header(slide6, "Supporting Agent: Gated Mitigation", "Incident Response")

    # Left: Orchestrator logic
    col1_box = slide6.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(5.6), Inches(4.5))
    tf6_col1 = col1_box.text_frame
    tf6_col1.word_wrap = True
    tf6_col1.margin_left = tf6_col1.margin_right = tf6_col1.margin_top = tf6_col1.margin_bottom = 0

    p_h6 = tf6_col1.paragraphs[0]
    p_h6.text = "Risk-Prioritized Response Planning"
    p_h6.font.size = Pt(20)
    p_h6.font.bold = True
    p_h6.font.color.rgb = TITLE_COLOR
    p_h6.space_after = Pt(15)

    bullets6 = [
        "Dynamic Risk Formula: risk = (anomaly score * 6.0) + (criticality weight), generating a standardized risk score scaling from 0 to 10.",
        "Asset Awareness: Integrates with an asset inventory database mapping device owners, IP subnets, and active CVE profiles.",
        "CVE Priority Override: Automatically bumps anomaly threshold parameters if targeting an asset with unpatched vulnerabilities."
    ]
    for bullet in bullets6:
        p_b = tf6_col1.add_paragraph()
        p_b.text = bullet
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_COLOR
        p_b.space_after = Pt(15)

    # Right: Response Actions Matrix
    card_r6 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(2.0), Inches(5.6), Inches(4.5))
    card_r6.fill.solid()
    card_r6.fill.fore_color.rgb = CARD_BG
    card_r6.line.color.rgb = CARD_BORDER
    card_r6.line.width = Pt(1.5)

    tb_m6 = slide6.shapes.add_textbox(Inches(6.98) + Inches(0.4), Inches(2.4), Inches(4.8), Inches(3.7))
    tf_m6 = tb_m6.text_frame
    tf_m6.word_wrap = True
    tf_m6.margin_left = tf_m6.margin_right = tf_m6.margin_top = tf_m6.margin_bottom = 0

    p_m6t = tf_m6.paragraphs[0]
    p_m6t.text = "Orchestrated Response Matrix"
    p_m6t.font.size = Pt(18)
    p_m6t.font.bold = True
    p_m6t.font.color.rgb = ACCENT_COLOR
    p_m6t.space_after = Pt(12)

    matrix = [
        ("Risk >= 8.5", "Isolate Endpoint", "Manual SOC Approval Required"),
        ("Risk >= 7.0", "Revoke Active Sessions", "Manual SOC Approval Required"),
        ("Risk >= 5.5", "Rate Limit IP Address", "Auto-Executed"),
        ("Risk >= 4.0", "Snapshot VM State", "Auto-Executed"),
        ("Risk < 4.0", "Log anomaly & flag alert", "Auto-Executed")
    ]
    for risk_range, action, status in matrix:
        p_row = tf_m6.add_paragraph()
        p_row.text = f"[{risk_range}] "
        p_row.font.size = Pt(12)
        p_row.font.bold = True
        p_row.font.color.rgb = TITLE_COLOR
        
        r_act = p_row.add_run()
        r_act.text = f"{action}  "
        r_act.font.bold = True
        r_act.font.color.rgb = TEXT_COLOR
        
        r_stat = p_row.add_run()
        r_stat.text = f"({status})"
        r_stat.font.italic = True
        r_stat.font.color.rgb = RGBColor(34, 197, 94) if "Auto" in status else RGBColor(239, 68, 68)
        p_row.space_after = Pt(8)


    # ================= Slide 7: SOC Operations Console =================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide7)
    add_slide_header(slide7, "SOC Analyst Threat Console", "Operations Control")

    # Left: Bullet list of features
    col1_box = slide7.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(5.6), Inches(4.5))
    tf7_col1 = col1_box.text_frame
    tf7_col1.word_wrap = True
    tf7_col1.margin_left = tf7_col1.margin_right = tf7_col1.margin_top = tf7_col1.margin_bottom = 0

    p_h7 = tf7_col1.paragraphs[0]
    p_h7.text = "Streamlit Operations Dashboard"
    p_h7.font.size = Pt(20)
    p_h7.font.bold = True
    p_h7.font.color.rgb = TITLE_COLOR
    p_h7.space_after = Pt(15)

    dashboard_bullets = [
        "Live Threat Feed: Real-time table viewing connection alerts, timestamps, severity scores, and attributed techniques.",
        "Interactive Timeline: Interactive scatter plot showing threat scoring history to monitor temporal surges.",
        "Gated Approval Queue: Clear card console allowing SOC analysts to approve, override, or dismiss pending high blast-radius actions.",
        "MITRE Vector Sandbox: Search playground enabling analysts to input raw text queries to test ChromaDB semantic resolution."
    ]
    for bullet in dashboard_bullets:
        p_b = tf7_col1.add_paragraph()
        p_b.text = bullet
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_COLOR
        p_b.space_after = Pt(12)

    # Right: Add space for dashboard screenshots in final packaging
    card_r7 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(2.0), Inches(5.6), Inches(4.5))
    card_r7.fill.solid()
    card_r7.fill.fore_color.rgb = CARD_BG
    card_r7.line.color.rgb = CARD_BORDER
    card_r7.line.width = Pt(1.5)

    tb_screen = slide7.shapes.add_textbox(Inches(6.98) + Inches(0.4), Inches(3.0), Inches(4.8), Inches(2.5))
    tf_screen = tb_screen.text_frame
    tf_screen.word_wrap = True
    tf_screen.margin_left = tf_screen.margin_right = tf_screen.margin_top = tf_screen.margin_bottom = 0

    p_sc1 = tf_screen.paragraphs[0]
    p_sc1.text = "💻 Interactive SOC Terminal"
    p_sc1.font.size = Pt(22)
    p_sc1.font.bold = True
    p_sc1.font.color.rgb = ACCENT_COLOR
    p_sc1.alignment = PP_ALIGN.CENTER
    p_sc1.space_after = Pt(15)

    p_sc2 = tf_screen.add_paragraph()
    p_sc2.text = "Visual dashboard running stream-lit on port 8501. Directly queries the local SQLite database for threat data logs."
    p_sc2.font.size = Pt(14)
    p_sc2.font.color.rgb = TEXT_COLOR
    p_sc2.alignment = PP_ALIGN.CENTER


    # ================= Slide 8: Quantitative Evaluation & Results =================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide8)
    add_slide_header(slide8, "Quantitative Evaluation Results", "Validation Metrics")

    # Left: Evaluation description
    col1_box = slide8.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(5.6), Inches(4.5))
    tf8_col1 = col1_box.text_frame
    tf8_col1.word_wrap = True
    tf8_col1.margin_left = tf8_col1.margin_right = tf8_col1.margin_top = tf8_col1.margin_bottom = 0

    p_h8 = tf8_col1.paragraphs[0]
    p_h8.text = "Robust Model Performance"
    p_h8.font.size = Pt(20)
    p_h8.font.bold = True
    p_h8.font.color.rgb = TITLE_COLOR
    p_h8.space_after = Pt(15)

    eval_bullets = [
        "Tested on Holdout Dataset: Evaluated on test holdout set (~22,000 instances) to reflect real deployment distributions.",
        "Local Semantic Matching: Obtained 65% correct classification matches on hard-constructed scenarios using vector matching alone.",
        "Volumetric Scan Success: 99.59% detection rate of probe/scanning attacks, and 88.13% of DoS traffic flows blocked."
    ]
    for bullet in eval_bullets:
        p_b = tf8_col1.add_paragraph()
        p_b.text = bullet
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_COLOR
        p_b.space_after = Pt(15)

    # Right: Results Data Card
    card_r8 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(2.0), Inches(5.6), Inches(4.5))
    card_r8.fill.solid()
    card_r8.fill.fore_color.rgb = CARD_BG
    card_r8.line.color.rgb = CARD_BORDER
    card_r8.line.width = Pt(1.5)

    tb_results = slide8.shapes.add_textbox(Inches(6.98) + Inches(0.4), Inches(2.4), Inches(4.8), Inches(3.7))
    tf_res = tb_results.text_frame
    tf_res.word_wrap = True
    tf_res.margin_left = tf_res.margin_right = tf_res.margin_top = tf_res.margin_bottom = 0

    p_rest = tf_res.paragraphs[0]
    p_rest.text = "Performance Metrics Summary"
    p_rest.font.size = Pt(18)
    p_rest.font.bold = True
    p_rest.font.color.rgb = ACCENT_COLOR
    p_rest.space_after = Pt(15)

    metrics_list = [
        ("• Normal Traffic Isolation Rate", "95.71%"),
        ("• Probe Attack Detection Rate", "99.59%"),
        ("• DoS Attack Detection Rate", "88.13%"),
        ("• Overall Anomaly F1-Score", "0.8381"),
        ("• Semantic Attribution Accuracy", "65.00%")
    ]
    for label, val in metrics_list:
        p_line = tf_res.add_paragraph()
        p_line.text = f"{label}: "
        p_line.font.size = Pt(14)
        p_line.font.color.rgb = TEXT_COLOR
        
        r_val = p_line.add_run()
        r_val.text = val
        r_val.font.bold = True
        r_val.font.color.rgb = TITLE_COLOR
        p_line.space_after = Pt(12)


    # ================= Slide 9: Why SentinelMind AI Wins =================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide9)
    add_slide_header(slide9, "Platform Architecture Highlights", "SentinelMind Edge")

    # Four small grid blocks for highlights
    card_w = Inches(5.6)
    card_h = Inches(2.1)
    
    highlights = [
        {
            "title": "⚡ Zero-Day Identification",
            "desc": "Detects novel threats without signatures by modeling normal network behaviors using unsupervised Isolation Forests.",
            "col": 0, "row": 0
        },
        {
            "title": "🔍 Context-Rich Alerts",
            "desc": "Enriches every single alert with named MITRE ATT&CK techniques, detailed prose descriptions, and full decision trails.",
            "col": 1, "row": 0
        },
        {
            "title": "🔒 Human-in-the-Loop Safety",
            "desc": "High blast-radius playbooks are automatically gated, letting analysts review metrics before executing network isolation.",
            "col": 0, "row": 1
        },
        {
            "title": "🛠️ Modular Technical Stack",
            "desc": "Uses light FastAPI services, persistent ChromaDB collections, sqlite3 storage, and an elegant Streamlit frontend.",
            "col": 1, "row": 1
        }
    ]

    for hl in highlights:
        x = Inches(0.75) if hl["col"] == 0 else Inches(6.98)
        y = Inches(2.0) if hl["row"] == 0 else Inches(4.4)
        
        card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.5)

        tb = slide9.shapes.add_textbox(x + Inches(0.25), y + Inches(0.25), card_w - Inches(0.5), card_h - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        p_t = tf.paragraphs[0]
        p_t.text = hl["title"]
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_COLOR
        p_t.space_after = Pt(8)

        p_d = tf.add_paragraph()
        p_d.text = hl["desc"]
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_COLOR


    # ================= Slide 10: Conclusion & Next Steps =================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide10)
    add_slide_header(slide10, "Hackathon Submission Summary", "Next Steps")

    # Left: Conclusion
    col1_box = slide10.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(5.6), Inches(4.5))
    tf10_col1 = col1_box.text_frame
    tf10_col1.word_wrap = True
    tf10_col1.margin_left = tf10_col1.margin_right = tf10_col1.margin_top = tf10_col1.margin_bottom = 0

    p_h10 = tf10_col1.paragraphs[0]
    p_h10.text = "Submission-Ready Package"
    p_h10.font.size = Pt(20)
    p_h10.font.bold = True
    p_h10.font.color.rgb = TITLE_COLOR
    p_h10.space_after = Pt(15)

    bullets10 = [
        "Fully Operational Codebase: Complete microservices, schema-validated alert pipeline, and live Streamlit interface.",
        "Standardized Benchmarks: Evaluation results documented on both ML performance and threat attribution profiles.",
        "Documentation Complete: In-depth setup instructions, Mermaid architecture files, and design limitations recorded."
    ]
    for bullet in bullets10:
        p_b = tf10_col1.add_paragraph()
        p_b.text = bullet
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_COLOR
        p_b.space_after = Pt(15)

    # Right: Roadmap
    card_r10 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(2.0), Inches(5.6), Inches(4.5))
    card_r10.fill.solid()
    card_r10.fill.fore_color.rgb = CARD_BG
    card_r10.line.color.rgb = CARD_BORDER
    card_r10.line.width = Pt(1.5)

    tb_roadmap = slide10.shapes.add_textbox(Inches(6.98) + Inches(0.4), Inches(2.4), Inches(4.8), Inches(3.7))
    tf_road = tb_roadmap.text_frame
    tf_road.word_wrap = True
    tf_road.margin_left = tf_road.margin_right = tf_road.margin_top = tf_road.margin_bottom = 0

    p_roadt = tf_road.paragraphs[0]
    p_roadt.text = "Platform Future Roadmap"
    p_roadt.font.size = Pt(18)
    p_roadt.font.bold = True
    p_roadt.font.color.rgb = ACCENT_COLOR
    p_roadt.space_after = Pt(15)

    roadmap_items = [
        ("• Production SIEM Integration", "Connecting ingest pipelines directly to active enterprise agents (Splunk, Elastic)."),
        ("• Reinforcement Learning Mitigations", "Training secondary agents to tune playbook blast-radius decisions based on human feedback loops."),
        ("• Multi-Protocol Pipeline Scaling", "Scaling parsing capabilities to support DNS, HTTP/HTTPS payload inspections, and raw PCAP captures.")
    ]
    for r_title, r_desc in roadmap_items:
        p_rt = tf_road.add_paragraph()
        p_rt.text = r_title
        p_rt.font.size = Pt(14)
        p_rt.font.bold = True
        p_rt.font.color.rgb = TEXT_COLOR
        
        p_rd = tf_road.add_paragraph()
        p_rd.text = r_desc
        p_rd.font.size = Pt(12)
        p_rd.font.color.rgb = MUTED_COLOR
        p_rd.space_after = Pt(12)

    # Save to disk
    output_path = os.path.join(os.path.dirname(__file__), "..", "..", "PITCH_DECK.pptx")
    output_path = os.path.abspath(output_path)
    prs.save(output_path)
    print(f"PowerPoint Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_pitch_deck()
